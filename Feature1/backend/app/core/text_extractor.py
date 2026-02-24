import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from typing import Optional

from app.models.file_metadata import FileMetadata
from app.utils.helpers import (
    calculate_sha256,
    extract_file_metadata,
    is_supported_file,
)
from app.utils.logger import logger


class TextExtractor:

    def extract_text(self, file_path: str) -> Optional[str]:
        ext = os.path.splitext(file_path)[1].lower()

        try:
            if ext == ".pdf":
                return self._extract_pdf(file_path)
            elif ext == ".docx":
                return self._extract_docx(file_path)
            elif ext == ".pptx":
                return self._extract_pptx(file_path)
            elif ext == ".txt":
                return self._extract_txt(file_path)
            else:
                return None

        except Exception as e:
            logger.error(f"Text extraction failed for {file_path}: {e}")
            return None

    def process_file(self, file_path: str) -> Optional[FileMetadata]:

        if not os.path.exists(file_path):
            logger.warning(f"File does not exist: {file_path}")
            return None

        if not is_supported_file(file_path):
            logger.info(f"Unsupported file skipped: {file_path}")
            return None

        try:
            metadata = extract_file_metadata(file_path)
            sha256_hash = calculate_sha256(file_path)
            extracted_text = self.extract_text(file_path)

            return FileMetadata(
                file_path=file_path,
                file_name=os.path.basename(file_path),
                file_extension=os.path.splitext(file_path)[1].lower(),
                size_bytes=metadata["size_bytes"],
                created_at=metadata["created_at"],
                modified_at=metadata["modified_at"],
                sha256_hash=sha256_hash,
                extracted_text=extracted_text,
            )

        except Exception as e:
            logger.error(f"Processing failed for {file_path}: {e}")
            return None

    def _extract_pdf(self, file_path: str) -> str:
        text = ""
        with fitz.open(file_path) as doc:
            for page in doc:
                text += page.get_text()
        return text.strip()

    def _extract_docx(self, file_path: str) -> str:
        doc = docx.Document(file_path)
        return "\n".join([para.text for para in doc.paragraphs]).strip()

    def _extract_pptx(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs).strip()

    def _extract_txt(self, file_path: str) -> str:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read().strip()